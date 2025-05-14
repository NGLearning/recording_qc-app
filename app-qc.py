
import streamlit as st
import whisper
import openai
from pptx import Presentation
from tempfile import NamedTemporaryFile
import subprocess
import os
import io
from openai import AzureOpenAI

# Azure setup (use secrets in Streamlit Cloud or local .streamlit/secrets.toml)
AZURE_API_KEY = st.secrets["AZURE_API_KEY"]
AZURE_ENDPOINT = st.secrets["AZURE_ENDPOINT"]
AZURE_DEPLOYMENT = st.secrets["AZURE_DEPLOYMENT"]

openai.api_key = AZURE_API_KEY
openai.api_base = AZURE_ENDPOINT
openai.api_type = "azure"
openai.api_version = "2025-03-01-preview"

client = AzureOpenAI(
  azure_endpoint = AZURE_ENDPOINT,
  api_key=AZURE_API_KEY,
  api_version="2025-03-01-preview"
)
# Convert MP4 to WAV
def convert_mp4_to_wav(mp4_bytes, output_path):
    with NamedTemporaryFile(delete=False, suffix=".mp4") as temp_file:
        temp_file.write(mp4_bytes)
        temp_file.flush()
        subprocess.call(['ffmpeg', '-i', temp_file.name, output_path])
    os.remove(temp_file.name)

# Transcribe Audio
def transcribe_audio(wav_path):
    model = whisper.load_model("base")
    result = model.transcribe(wav_path)
    return result['text']

# Extract Text from PPTX
def extract_text_from_pptx(pptx_file):
    prs = Presentation(pptx_file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# Call LLM for QC
def quality_check(video_transcript, ppt_text):
    user_input = f"""
The video content is based on the following presentation slides:
--- PPT Content Start ---
{ppt_text}
--- PPT Content End ---

Below is the transcript of the video that needs to be quality checked:
--- Transcript Start ---
{video_transcript}
--- Transcript End ---

You are part of a data science team at an EdTech company. Your task is to review the above video transcript and its corresponding PPT content for quality, using the checklist below.

### Checklist:

**General Checks:**
1. Identify the topics and sub-topics taught in the session, in order and in the form of a hierarchical list
2. Interpret what the learning outcomes were from the session (use Bloom’s Taxonomy to synthesize outcomes)
3. Identify any logistic / technical issues that occurred in the recording
4. Identify any comments that are potentially racist, inflammatory, sexually inappropriate, biased, or polarizing
5. Assess whether there is any mismatch in the faculty's tone, particularly signs of low energy or engagement

**For Conceptual Teaching:**
6. Scrutinize for any flow-related inaccuracies or misleading information presented by the mentor
7. Determine the pace of the recording by looking at how quickly the topics were progressing in the teaching

**For Code Walkthrough:**
Note: The following items apply **only if the video is a code walkthrough**. If this is a conceptual explanation and no code walkthrough is present, mark each item as "N/A - Not a code walkthrough video".
8. Is the walkthrough of the code sequential and line-by-line? Does the faculty skip explaining any sections?
9. Are all essential features of the code discussed? Is there a comprehensive explanation of key functions, methods, or classes?
10. Is the faculty explaining all the logic present in the code clearly?
11. Does the faculty talk about the functions being used, why they were written and how they work?
12. When using in-built functions and methods, does the faculty explain what each parameter does and why we need them?
13. Is there a clear explanation of the expected output from the code? Does the instructor validate output results and explain discrepancies, if any?
14. Are real-world examples or use cases provided to contextualize the code?

### Output Format:
- For each checklist item, respond with ✅ or ❌.
- If ❌, add 1–2 lines explaining the issue.
- At the end, provide a short bullet list of what went wrong and how to improve in two separate heading.
- Do not use special characters like ### or ***.

Generate only the checklist and bullet point feedback. Do not include unnecessary explanation or summary.
"""

    response = client.chat.completions.create(
        model="gpt-4o",  # or "gpt-4o-mini" if needed
        messages=[
            {"role": "system", "content": "You are a helpful assistant."},
            {"role": "user", "content": user_input}
        ],
        max_tokens=4096,
        temperature=0.5
    )

    return response.choices[0].message.content.strip()

# Streamlit UI
st.title("🎧 Recording Quality Check with LLM")

mp4_file = st.file_uploader("Upload MP4 Video", type=["mp4"])
pptx_file = st.file_uploader("Upload PPTX Slides", type=["pptx"])

if st.button("Run Quality Check") and mp4_file and pptx_file:
    with st.spinner("Processing..."):
        convert_mp4_to_wav(mp4_file.read(), "audio.wav")
        transcript = transcribe_audio("audio.wav")
        ppt_text = extract_text_from_pptx(pptx_file)
        report = quality_check(transcript, ppt_text)

        st.success("✅ Quality Check Complete")
        st.text_area("QC Report", report, height=400)

        # Enable download

        st.download_button(
            label="📥 Download QC Report",
            data=report,
            file_name="QC_Report.txt",
            mime="text/plain"
        )        
